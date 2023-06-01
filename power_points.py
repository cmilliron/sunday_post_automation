from pptx import Presentation

file_path = r"Resources\powerpoints\Psalm 031 (1-16).pptx"
verse_set = [str(i) for i in range(1, 100, 1)]


# save_file = "hold.pptx"
# p = slides.Presentation(file_path)
# text_runs will be populated with a list of strings,
# one for each text run in presentation


def slides_to_test(presentation):
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.bold:
                        # Print the bold text to the console
                        print("{y}" + f"{run.text.strip()}" + "{/y}")
                        text_runs.append("{y}" + f"{run.text.strip()}" + "{/y}")
                    else:
                        # Print the normal text to the console
                        # print(run.text.strip())
                        text_runs.append(run.text.strip())
    return text_runs


def text_to_openlp(text, bulletin=False):
    line = ""
    intro_text = True
    new_line = False
    for t in text:
        if t in verse_set:
            intro_text = False
            # print(line)
            # print("[===]")
            # line += t
        if intro_text == False:
            if t == "":
                line += "\n[===]\n"
                new_line = True
            else:
                if (t[-1] in verse_set) or (t[-1] == "R"):
                    output = t[0:-2]
                else:
                    output = t
                output.strip()
                if line[:-6] == "[===]\n":
                    line += output.rstrip("\n")
                else:
                    if new_line == True:
                        line = line + output.rstrip("\n")
                        new_line = False
                    else:
                        line = line + " " + output.rstrip("\n")
    return line


def text_to_openlp_psalm(text):
    line = ""
    for t in text:
        if t[0] in verse_set and t[1] == '.':
            print(line)
            print("[===]")
            line = t[3:]
        elif line == "":
            line = t
        else:
            t.strip()
            line = line + " " + t
    return line


if __name__ == "__main__":
    prs = Presentation(file_path)
    prs_text = slides_to_test(prs)
    print(text_to_openlp(prs_text))
    print(text_to_openlp(prs_text, bulletin=True))
    # print(prs_text)
